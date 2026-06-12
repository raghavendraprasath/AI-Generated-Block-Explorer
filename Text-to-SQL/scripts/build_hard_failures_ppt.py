#!/usr/bin/env python3
"""Build hard_failures.pptx + per-case screenshots for class presentation."""

from __future__ import annotations

import json
import sqlite3
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
HARD_PATH = ROOT / "tests" / "hard_failures.json"
SLIDES_DIR = ROOT / "slides"
IMG_DIR = SLIDES_DIR / "ppt_assets"
PPT_PATH = SLIDES_DIR / "hard_failures.pptx"
DEFAULT_DB = Path.home() / "hw3-data" / "blockchain.db"

BG = (13, 17, 23)
PANEL = (22, 27, 34)
GREEN = (63, 185, 80)
RED = (248, 81, 73)
ORANGE = (247, 147, 26)
TEXT = (230, 237, 243)
MUTED = (139, 148, 158)

RGB_BG = RGBColor(*BG)
RGB_PANEL = RGBColor(*PANEL)
RGB_TEXT = RGBColor(*TEXT)
RGB_MUTED = RGBColor(*MUTED)
RGB_ORANGE = RGBColor(*ORANGE)
def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSansMono-Bold.ttf"
        if bold
        else "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationMono-Bold.ttf"
        if bold
        else "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf",
        "C:/Windows/Fonts/consola.ttf",
        "C:/Windows/Fonts/consolab.ttf" if bold else "C:/Windows/Fonts/consola.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def _wrap(text: str, width: int = 92) -> list[str]:
    lines: list[str] = []
    for raw in text.splitlines():
        if not raw.strip():
            lines.append("")
            continue
        lines.extend(textwrap.wrap(raw, width=width) or [""])
    return lines


def render_terminal_image(
    path: Path,
    *,
    title: str,
    correct_sql: str,
    correct_answer: str,
    wrong_sql: str,
    wrong_answer: str,
) -> None:
    font = _font(18)
    font_bold = _font(20, bold=True)
    font_title = _font(24, bold=True)

    lines: list[tuple[str, tuple[int, int, int]]] = []
    lines.append((title, ORANGE))
    lines.append(("", TEXT))
    lines.append(("✓ CORRECT SQL", GREEN))
    for line in _wrap(correct_sql, 88):
        lines.append((f"  {line}", TEXT))
    lines.append((f"  → {correct_answer}", GREEN))
    lines.append(("", TEXT))
    lines.append(("✗ LLM SQL (wrong)", RED))
    for line in _wrap(wrong_sql, 88):
        lines.append((f"  {line}", TEXT))
    lines.append((f"  → {wrong_answer}", RED))

    line_h = 28
    pad = 28
    width = 1280
    height = pad * 2 + line_h * len(lines) + 20
    img = Image.new("RGB", (width, height), BG)
    draw = ImageDraw.Draw(img)

    y = pad
    for i, (line, color) in enumerate(lines):
        f = font_title if i == 0 else font_bold if line.startswith(("✓", "✗")) else font
        draw.text((pad, y), line, fill=color, font=f)
        y += line_h

    # Subtle border so screenshots read clearly on dark slides
    draw.rectangle(
        [(2, 2), (width - 3, height - 3)],
        outline=ORANGE,
        width=2,
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path)


def _new_dark_slide(prs: Presentation):
    """Blank slide with GitHub-dark background and orange top accent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGB_BG

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGB_ORANGE
    accent.line.fill.background()
    return slide


def _add_panel(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGB_PANEL
    panel.line.color.rgb = RGBColor(48, 54, 61)
    panel.line.width = Pt(1)


CASE_META = [
    ("Case 1 — Column selection", "case1_column_selection.png"),
    ("Case 2 — Aggregation logic", "case2_aggregation.png"),
    ("Case 3 — Domain enums", "case3_script_types.png"),
]


def load_cases(db_path: Path) -> list[dict]:
    cases = json.loads(HARD_PATH.read_text(encoding="utf-8"))
    if not db_path.exists():
        return cases
    conn = sqlite3.connect(db_path)
    try:
        for case in cases:
            case["expected_answer"] = conn.execute(case["expected_sql"]).fetchone()[0]
            case["incorrect_answer"] = conn.execute(case["incorrect_sql"]).fetchone()[0]
    finally:
        conn.close()
    return cases


def build_images(cases: list[dict]) -> list[Path]:
    paths: list[Path] = []
    for (title, filename), case in zip(CASE_META, cases, strict=True):
        out = IMG_DIR / filename
        render_terminal_image(
            out,
            title=title,
            correct_sql=case["expected_sql"],
            correct_answer=str(case["expected_answer"]),
            wrong_sql=case["incorrect_sql"],
            wrong_answer=str(case["incorrect_answer"]),
        )
        paths.append(out)
    return paths


def _add_title_slide(prs: Presentation) -> None:
    slide = _new_dark_slide(prs)
    _add_panel(slide, 0.55, 1.15, 8.9, 4.6)

    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(8.3), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "When Text-to-SQL Fails"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGB_TEXT

    for text, size, color in [
        ("INFO7500 · Homework 3 · Block Explorer AI", 22, RGB_ORANGE),
        (
            "Three hard cases where incorrect SQL still returns a plausible wrong answer",
            18,
            RGB_MUTED,
        ),
        ("Reference: Spider 2.0 Text-to-SQL leaderboard", 15, RGB_MUTED),
        ("Data: tests/hard_failures.json", 15, RGB_MUTED),
    ]:
        para = tf.add_paragraph()
        para.text = text
        para.font.size = Pt(size)
        para.font.color.rgb = color
        para.space_before = Pt(14)


def _add_case_slide(
    prs: Presentation,
    *,
    heading: str,
    question: str,
    why: str,
    image_path: Path,
) -> None:
    slide = _new_dark_slide(prs)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(9), Inches(0.55))
    tp = title.text_frame.paragraphs[0]
    tp.text = heading
    tp.font.size = Pt(30)
    tp.font.bold = True
    tp.font.color.rgb = RGB_ORANGE

    _add_panel(slide, 0.45, 0.78, 9.1, 0.72)
    qbox = slide.shapes.add_textbox(Inches(0.65), Inches(0.92), Inches(8.7), Inches(0.55))
    qp = qbox.text_frame.paragraphs[0]
    qp.text = f"Question: {question}"
    qp.font.size = Pt(17)
    qp.font.bold = True
    qp.font.color.rgb = RGB_TEXT

    slide.shapes.add_picture(str(image_path), Inches(0.35), Inches(1.62), width=Inches(9.3))

    _add_panel(slide, 0.45, 5.05, 9.1, 0.95)
    wbox = slide.shapes.add_textbox(Inches(0.65), Inches(5.2), Inches(8.7), Inches(0.75))
    wp = wbox.text_frame.paragraphs[0]
    wp.text = f"Why it's hard: {why}"
    wp.font.size = Pt(15)
    wp.font.color.rgb = RGB_MUTED
    wp.font.bold = False


def _add_takeaway_slide(prs: Presentation) -> None:
    slide = _new_dark_slide(prs)
    _add_panel(slide, 0.55, 0.85, 8.9, 5.5)

    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.1), Inches(8.3), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Takeaway"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGB_ORANGE

    bullets = [
        "Schema in context is not enough for free LLMs",
        "Failures: wrong column, wrong aggregation, wrong enum",
        "Wrong SQL can still execute and look plausible",
        "Golden SQL tests (12/12) validate the database; live LLM is the weak link",
        "Block Explorer AI · Text-to-SQL/",
    ]
    for bullet in bullets:
        para = tf.add_paragraph()
        para.text = f"• {bullet}"
        para.level = 0
        para.font.size = Pt(21)
        para.font.color.rgb = RGB_TEXT
        para.space_before = Pt(12)


def build_ppt(cases: list[dict], image_paths: list[Path], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs)
    for (heading, _), case, image_path in zip(CASE_META, cases, image_paths, strict=True):
        _add_case_slide(
            prs,
            heading=heading,
            question=case["question"],
            why=case["why_hard"],
            image_path=image_path,
        )
    _add_takeaway_slide(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> int:
    import sys

    db = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_DB
    cases = load_cases(db)
    images = build_images(cases)
    build_ppt(cases, images, PPT_PATH)
    print(f"Wrote {PPT_PATH}")
    for img in images:
        print(f"  screenshot: {img}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
